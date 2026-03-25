#!/usr/bin/env python3

import sys
# Add user site-packages for Apache service (which runs as different user)
sys.path.insert(0, r'C:\Users\RAGHIB\AppData\Roaming\Python\Python314\site-packages')

import argparse
import os
import shutil
import subprocess
import tempfile
from dataclasses import dataclass
from pathlib import Path


def _fail(msg: str, code: int = 2):
    print(msg, file=sys.stderr)
    raise SystemExit(code)


def _require_file(p: str, label: str):
    if not p or not os.path.isfile(p):
        _fail(f"{label} not found: {p}")


def _require_dir(p: str, label: str):
    if not p or not os.path.isdir(p):
        _fail(f"{label} not found: {p}")


def _safe_int_list(values):
    out = []
    for v in values or []:
        try:
            n = int(str(v).strip())
        except Exception:
            continue
        if 1 <= n <= 9:
            out.append(n)
    # keep order but dedupe
    seen = set()
    uniq = []
    for n in out:
        if n in seen:
            continue
        seen.add(n)
        uniq.append(n)
    return uniq


@dataclass
class TargetShape:
    name: str
    left: int
    top: int
    width: int
    height: int


def _walk_shapes(shapes):
    """Yield shapes recursively, including group shapes' children."""
    for shape in shapes:
        yield shape
        try:
            # GroupShape has .shapes
            sub = getattr(shape, "shapes", None)
            if sub is not None:
                for child in _walk_shapes(sub):
                    yield child
        except Exception:
            continue


def _collect_named_shapes(prs, shape_name: str):
    found = []
    for slide in prs.slides:
        for shape in _walk_shapes(slide.shapes):
            try:
                if getattr(shape, "name", None) == shape_name:
                    found.append(shape)
            except Exception:
                continue
    return found


def _set_text(shape, text: str, template_type: str = None):
    """Set text on shape with specific font styling for each template type."""
    try:
        tf = shape.text_frame
    except Exception:
        return
    
    # Clear and set new text
    try:
        tf.clear()
    except Exception:
        pass

    try:
        tf.text = text
    except Exception:
        # fallback
        try:
            p = tf.paragraphs[0]
            p.text = text
        except Exception:
            return
    
    # Apply template-specific font styling
    try:
        from pptx.util import Pt
        
        p = tf.paragraphs[0]
        if p.runs:
            font = p.runs[0].font
        else:
            font = p.font
        
        # Always center align
        p.alignment = 2  # PP_ALIGN_CENTER = 2 (center)
        
        # Also set vertical centering for text frame
        try:
            tf.margin_top = 0
            tf.margin_bottom = 0
            tf.margin_left = 0
            tf.margin_right = 0
            tf.vertical_anchor = 1  # MSO_ANCHOR_MIDDLE = 1 (center)
        except Exception:
            pass
        
        if template_type == 'QR':
            # QR: Bahnschrift SemiCondensed / 16 / Bold / Text shadow / centered
            font.name = 'Bahnschrift SemiCondensed'
            font.size = Pt(16)
            font.bold = True
        elif template_type == 'TAG':
            # Tag: Bahnschrift SemiCondensed / 12 / Bold / Text shadow / centered
            font.name = 'Bahnschrift SemiCondensed'
            font.size = Pt(12)
            font.bold = True
        else:
            # Default fallback
            font.bold = True
            p.alignment = 1  # Centered
        
        # Add text shadow effect (glow effect in PowerPoint)
        try:
            from pptx.dml.color import RGBColor
            from pptx.enum.dml import MSO_THEME_COLOR_INDEX
            
            # Create a subtle shadow effect
            shadow = font.shadow
            shadow.inherit = False
            shadow.visible = True
            shadow.blur_radius = Pt(3)
            shadow.distance = Pt(2)
            shadow.angle = 45
            shadow.color.rgb = RGBColor(128, 128, 128)  # Gray shadow
        except Exception:
            # If shadow creation fails, continue without it
            pass
            
    except Exception:
        pass


def _replace_with_picture(shape, image_path: str, preserve_aspect: bool = True, match_height: bool = False):
    """Replace a shape by adding a picture at the same position.
    
    If match_height is True, scale based on height only (matching the shape height exactly).
    If preserve_aspect is True, scale to fit within bounds while preserving aspect ratio.
    Otherwise stretch to fit.
    """
    try:
        from PIL import Image
    except Exception:
        preserve_aspect = False
        match_height = False
    
    slide = shape.part.slide
    left = shape.left
    top = shape.top
    width = shape.width
    height = shape.height
    
    if match_height or preserve_aspect:
        try:
            with Image.open(image_path) as img:
                img_w, img_h = img.size
                if img_w > 0 and img_h > 0:
                    if match_height:
                        # Scale to match height exactly, keep aspect ratio
                        scale = height / img_h
                        new_h = int(img_h * scale)
                        new_w = int(img_w * scale)
                    else:
                        # Fit within bounds while preserving aspect ratio
                        scale_w = width / img_w
                        scale_h = height / img_h
                        scale = min(scale_w, scale_h)
                        new_w = int(img_w * scale)
                        new_h = int(img_h * scale)
                    
                    # Center the image in the original bbox
                    new_left = left + (width - new_w) // 2
                    new_top = top + (height - new_h) // 2
                    slide.shapes.add_picture(image_path, new_left, new_top, width=new_w, height=new_h)
                    # remove original element
                    try:
                        el = shape._element
                        el.getparent().remove(el)
                    except Exception:
                        pass
                    return
        except Exception:
            pass
    
    # Fallback: original behavior (stretch to fit)
    slide.shapes.add_picture(image_path, left, top, width=width, height=height)
    
    # remove original element
    try:
        el = shape._element
        el.getparent().remove(el)
    except Exception:
        # If we can't remove, leave it (will overlap)
        pass


def _generate_qr_png(out_path: str, payload: str, center_logo_path: str = None):
    """Generate QR code matching the qr-code-styling configuration."""
    try:
        import qrcode
        from qrcode.image.styledpil import StyledPilImage
        from qrcode.image.styles.moduledrawers import SquareModuleDrawer
        from qrcode.image.styles.colormasks import SolidFillColorMask
        from qrcode.constants import ERROR_CORRECT_H
        from PIL import Image
    except Exception as e:
        _fail(f"Missing python dependency for QR generation. {e}")
    
    # Colors - matching the configuration
    DOT_COLOR = (0, 0, 0)          # #000000 - black dots
    CORNER_COLOR = (251, 157, 1)    # #FB9D01 - orange corners  
    BACKGROUND_COLOR = (255, 255, 255)  # White background (will be made transparent)
    
    # Create QR with high error correction for logo overlay
    qr = qrcode.QRCode(
        version=None,
        error_correction=ERROR_CORRECT_H,  # Best error correction (~30%)
        box_size=11,
        border=1,
    )
    qr.add_data(payload)
    qr.make(fit=True)
    
    # Generate styled QR with square dots and transparent background
    # Note: We use black for all dots, the corner colors are handled by the library's eye drawers
    color_mask = SolidFillColorMask(
        front_color=DOT_COLOR,
        back_color=BACKGROUND_COLOR
    )
    
    # Use square module drawer for all dots (matching "square" type)
    img = qr.make_image(
        image_factory=StyledPilImage,
        module_drawer=SquareModuleDrawer(),
        color_mask=color_mask
    )
    img = img.convert('RGBA')
    
    # Make white background transparent
    datas = img.getdata()
    newData = []
    for item in datas:
        if item[0] > 240 and item[1] > 240 and item[2] > 240:
            newData.append((255, 255, 255, 0))
        else:
            newData.append(item)
    img.putdata(newData)
    
    # Add center logo if provided
    if center_logo_path and os.path.isfile(center_logo_path):
        try:
            logo = Image.open(center_logo_path).convert("RGBA")
            
            # Calculate logo size (27% of QR size)
            qr_size = min(img.size)
            logo_size = int(qr_size * 0.27)
            logo = logo.resize((logo_size, logo_size), Image.LANCZOS)
            
            # Calculate center position
            center_x = (img.size[0] - logo_size) // 2
            center_y = (img.size[1] - logo_size) // 2
            
            # Create mask for clearing logo area (rounded rectangle)
            from PIL import ImageDraw
            mask = Image.new('L', img.size, 0)
            mask_draw = ImageDraw.Draw(mask)
            
            # Margin of 3px around logo
            margin = 3
            clear_box = [
                center_x - margin,
                center_y - margin,
                center_x + logo_size + margin,
                center_y + logo_size + margin
            ]
            mask_draw.rounded_rectangle(clear_box, radius=15, fill=255)
            
            # Clear the area
            img.paste((0, 0, 0, 0), (0, 0), mask)
            
            # Paste logo on top
            img.paste(logo, (center_x, center_y), logo)
            
        except Exception as e:
            print(f"Logo overlay failed: {e}", file=sys.stderr)
            pass
    
    # Save as PNG with transparency
    img.save(out_path, 'PNG')


def _compose_pictograms(out_path: str, pictograms_dir: str, ids: list[int], target_w: int = 1200, target_h: int = 600):
    try:
        from PIL import Image
    except Exception as e:
        _fail(f"Missing python dependency for image composition (Pillow). {e}")

    ids = _safe_int_list(ids)
    if not ids:
        # create transparent 1x1 to avoid crashing replace
        img = Image.new("RGBA", (1, 1), (0, 0, 0, 0))
        img.save(out_path)
        return

    pics = []
    for n in ids:
        p = os.path.join(pictograms_dir, f"{n}.png")
        if not os.path.isfile(p):
            _fail(f"Pictogram not found: {p}")
        pics.append(Image.open(p).convert("RGBA"))

    # Fit all pictograms into target box horizontally with spacing.
    spacing = max(10, int(target_w * 0.03))
    available_w = max(1, target_w - spacing * (len(pics) - 1))

    # each pictogram max width and height
    per_w = max(1, available_w // len(pics))
    per_h = max(1, target_h)

    resized = []
    for im in pics:
        w, h = im.size
        if w <= 0 or h <= 0:
            continue
        scale = min(per_w / w, per_h / h)
        nw = max(1, int(w * scale))
        nh = max(1, int(h * scale))
        resized.append(im.resize((nw, nh), Image.LANCZOS))

    total_w = sum(im.size[0] for im in resized) + spacing * (len(resized) - 1)
    max_h = max(im.size[1] for im in resized)

    canvas_w = min(target_w, max(target_w, total_w))
    canvas_h = min(target_h, max(target_h, max_h))

    canvas = Image.new("RGBA", (canvas_w, canvas_h), (0, 0, 0, 0))

    x = (canvas_w - total_w) // 2
    for im in resized:
        y = (canvas_h - im.size[1]) // 2
        canvas.alpha_composite(im, (x, y))
        x += im.size[0] + spacing

    canvas.save(out_path)


def _pptx_to_pdf(pptx_path: str, out_pdf_path: str, soffice_bin: str):
    out_pdf_path = os.path.abspath(out_pdf_path)
    out_dir = os.path.dirname(out_pdf_path)
    os.makedirs(out_dir, exist_ok=True)

    with tempfile.TemporaryDirectory(prefix="sds_pptx2pdf_") as tmpdir:
        tmpdir = os.path.abspath(tmpdir)

        cmd = [
            soffice_bin,
            "--headless",
            "--nologo",
            "--nolockcheck",
            "--norestore",
            "--convert-to",
            "pdf",
            "--outdir",
            tmpdir,
            pptx_path,
        ]

        proc = subprocess.run(cmd, capture_output=True, text=True)
        if proc.returncode != 0:
            _fail(f"LibreOffice conversion failed. stdout={proc.stdout} stderr={proc.stderr}")

        # libreoffice outputs with same base name
        base = os.path.splitext(os.path.basename(pptx_path))[0]
        produced = os.path.join(tmpdir, base + ".pdf")
        if not os.path.isfile(produced):
            # fallback: find any pdf
            pdfs = [p for p in os.listdir(tmpdir) if p.lower().endswith(".pdf")]
            if len(pdfs) == 1:
                produced = os.path.join(tmpdir, pdfs[0])
            else:
                _fail(f"LibreOffice did not produce expected PDF in {tmpdir}")

        shutil.copyfile(produced, out_pdf_path)


def _fill_qr_template(qr_template: str, product_name: str, qr_png: str, pictos_png: str, out_pptx: str):
    try:
        from pptx import Presentation
    except Exception as e:
        _fail(f"Missing python dependency for PPTX editing (python-pptx). {e}")

    prs = Presentation(qr_template)

    # QR Code
    qr_shapes = _collect_named_shapes(prs, "QR Code")
    if not qr_shapes:
        _fail('Shape named "QR Code" not found in QR template')
    for s in qr_shapes:
        _replace_with_picture(s, qr_png)

    # Product name
    name_shapes = _collect_named_shapes(prs, "Product name")
    if not name_shapes:
        _fail('Shape named "Product name" not found in QR template')
    for s in name_shapes:
        _set_text(s, product_name, template_type='QR')

    # Pictograms container
    pic_shapes = _collect_named_shapes(prs, "Pictograms")
    if not pic_shapes:
        _fail('Shape named "Pictograms" not found in QR template')
    for s in pic_shapes:
        _replace_with_picture(s, pictos_png, preserve_aspect=True)

    prs.save(out_pptx)


def _fill_tags_template(tag_template: str, product_name: str, qr_png: str, pictos_png: str, out_pptx: str):
    try:
        from pptx import Presentation
    except Exception as e:
        _fail(f"Missing python dependency for PPTX editing (python-pptx). {e}")

    prs = Presentation(tag_template)

    # Fill all repeated shapes across groups
    qr_shapes = _collect_named_shapes(prs, "QR Code")
    name_shapes = _collect_named_shapes(prs, "Product name")
    pic_shapes = _collect_named_shapes(prs, "Pictograms")

    if not qr_shapes:
        _fail('Shape named "QR Code" not found in Tags template')
    if not name_shapes:
        _fail('Shape named "Product name" not found in Tags template')
    if not pic_shapes:
        _fail('Shape named "Pictograms" not found in Tags template')

    for s in qr_shapes:
        _replace_with_picture(s, qr_png)

    for s in name_shapes:
        _set_text(s, product_name, template_type='TAG')

    for s in pic_shapes:
        _replace_with_picture(s, pictos_png, preserve_aspect=True)

    prs.save(out_pptx)


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--qr-template", required=True)
    ap.add_argument("--tag-template", required=True)
    ap.add_argument("--pictograms-dir", required=True)
    ap.add_argument("--product-name", required=True)
    ap.add_argument("--public-url", required=True)
    ap.add_argument("--out-qr-pdf", required=True)
    ap.add_argument("--out-tag-pdf", required=True)
    ap.add_argument("--doc-id", required=True)
    ap.add_argument("--picto", action="append", default=[])
    ap.add_argument("--soffice", default=os.environ.get("SDS_LIBREOFFICE_BIN", "soffice"))
    ap.add_argument("--center-logo", default=os.environ.get("SDS_CENTER_LOGO", ""))

    args = ap.parse_args()

    _require_file(args.qr_template, "QR template")
    _require_file(args.tag_template, "Tags template")
    _require_dir(args.pictograms_dir, "Pictograms dir")

    # Default center logo path if not provided
    center_logo_path = args.center_logo
    if not center_logo_path:
        # Try default location relative to script
        script_dir = os.path.dirname(os.path.abspath(__file__))
        default_logo = os.path.join(script_dir, "..", "storage", "app", "Templates", "Center logo.jpg")
        default_logo = os.path.abspath(default_logo)
        if os.path.isfile(default_logo):
            center_logo_path = default_logo

    pictos = _safe_int_list(args.picto)

    with tempfile.TemporaryDirectory(prefix=f"sds_gen_{args.doc_id}_") as tmpdir:
        tmpdir = os.path.abspath(tmpdir)

        qr_png = os.path.join(tmpdir, "qr.png")
        pictos_png = os.path.join(tmpdir, "pictos.png")

        _generate_qr_png(qr_png, args.public_url, center_logo_path if center_logo_path and os.path.isfile(center_logo_path) else None)
        # high-res composition; PowerPoint will scale down into bbox
        _compose_pictograms(pictos_png, args.pictograms_dir, pictos, 1920, 1248)

        qr_pptx_out = os.path.join(tmpdir, "qr_out.pptx")
        tag_pptx_out = os.path.join(tmpdir, "tag_out.pptx")

        _fill_qr_template(args.qr_template, args.product_name, qr_png, pictos_png, qr_pptx_out)
        _fill_tags_template(args.tag_template, args.product_name, qr_png, pictos_png, tag_pptx_out)

        _pptx_to_pdf(qr_pptx_out, args.out_qr_pdf, args.soffice)
        _pptx_to_pdf(tag_pptx_out, args.out_tag_pdf, args.soffice)

    print("ok")


if __name__ == "__main__":
    main()
