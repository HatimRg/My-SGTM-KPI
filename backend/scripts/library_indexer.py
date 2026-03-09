#!/usr/bin/env python3

import argparse
import json
import os
import re
import sys


def _detect_language(text: str) -> str:
    try:
        from langdetect import detect  # type: ignore

        lang = detect(text)
        if lang in {"en", "fr", "ar"}:
            return lang
        return lang
    except Exception:
        return "unknown"


def _extract_text(file_path: str, ext: str) -> str:
    ext = ext.lower().lstrip(".")

    if ext == "txt":
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()

    if ext == "pdf":
        try:
            from pdfminer.high_level import extract_text  # type: ignore

            return extract_text(file_path) or ""
        except Exception as e:
            raise RuntimeError(
                "PDF extraction requires pdfminer.six. Install: pip install pdfminer.six"
            ) from e

    if ext == "docx":
        try:
            import docx  # type: ignore

            d = docx.Document(file_path)
            parts = []
            for p in d.paragraphs:
                if p.text:
                    parts.append(p.text)
            return "\n".join(parts)
        except Exception as e:
            raise RuntimeError(
                "DOCX extraction requires python-docx. Install: pip install python-docx"
            ) from e

    if ext == "xlsx":
        try:
            import openpyxl  # type: ignore

            wb = openpyxl.load_workbook(file_path, data_only=True)
            parts = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    for cell in row:
                        if cell is None:
                            continue
                        s = str(cell).strip()
                        if s:
                            parts.append(s)
            return "\n".join(parts)
        except Exception as e:
            raise RuntimeError(
                "XLSX extraction requires openpyxl. Install: pip install openpyxl"
            ) from e

    if ext == "pptx":
        try:
            from pptx import Presentation  # type: ignore

            prs = Presentation(file_path)
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        t = str(getattr(shape, "text") or "").strip()
                        if t:
                            parts.append(t)
            return "\n".join(parts)
        except Exception as e:
            raise RuntimeError(
                "PPTX extraction requires python-pptx. Install: pip install python-pptx"
            ) from e

    raise RuntimeError(f"Unsupported extension for indexing: {ext}")


def _clean_text(text: str) -> str:
    text = re.sub(r"\s+", " ", text or "").strip()
    return text


def _extract_keywords(text: str, lang: str, top_n: int) -> list:
    try:
        from keybert import KeyBERT  # type: ignore

        model_name = os.environ.get("LIBRARY_KEYBERT_MODEL", "all-MiniLM-L6-v2")
        kw_model = KeyBERT(model=model_name)

        kws = kw_model.extract_keywords(
            text,
            top_n=top_n,
            stop_words="english" if lang == "en" else None,
        )
        out = []
        for kw, score in kws:
            kw = str(kw).strip()
            if kw:
                out.append(kw)
        return out
    except Exception as e:
        raise RuntimeError(
            "Keyword extraction requires KeyBERT + its dependencies. Install: pip install keybert sentence-transformers"
        ) from e


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--file", required=True)
    parser.add_argument("--ext", required=True)
    parser.add_argument("--top", type=int, default=12)
    args = parser.parse_args()

    file_path = args.file
    ext = args.ext
    top_n = args.top

    if not os.path.isfile(file_path):
        print(json.dumps({"ok": False, "error": "File not found"}))
        return 2

    try:
        raw_text = _extract_text(file_path, ext)
        text = _clean_text(raw_text)
        if not text:
            print(json.dumps({"ok": False, "error": "No text extracted"}))
            return 3

        lang = _detect_language(text)
        keywords = _extract_keywords(text, lang, top_n)

        print(
            json.dumps(
                {
                    "ok": True,
                    "language": lang,
                    "keywords": keywords,
                    "text_length": len(text),
                },
                ensure_ascii=False,
            )
        )
        return 0
    except Exception as e:
        print(json.dumps({"ok": False, "error": str(e)}), file=sys.stdout)
        return 1


if __name__ == "__main__":
    raise SystemExit(main())
